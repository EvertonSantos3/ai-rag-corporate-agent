"""
Extração de conteúdo por formato de arquivo (etapa 2 do desafio).

Cada loader recebe um caminho de arquivo e devolve texto limpo, pronto
para chunking. Mantemos um loader por formato para isolar dependências
(ex: só quem processa PDF precisa do pypdf instalado).
"""
from __future__ import annotations

import csv
import json
from pathlib import Path


class LoaderNaoSuportadoError(Exception):
    pass


def carregar_markdown(caminho: Path) -> str:
    """Markdown já é texto legível; removemos apenas marcações de sintaxe pesada."""
    texto = caminho.read_text(encoding="utf-8")
    return texto.strip()


def carregar_csv(caminho: Path) -> str:
    """Converte CSV em texto linha a linha, repetindo os cabeçalhos (facilita busca semântica)."""
    linhas_texto = []
    with caminho.open(newline="", encoding="utf-8") as f:
        leitor = csv.DictReader(f)
        for linha in leitor:
            partes = [f"{campo}: {valor}" for campo, valor in linha.items()]
            linhas_texto.append("; ".join(partes))
    return "\n".join(linhas_texto)


def carregar_json(caminho: Path) -> str:
    """Converte um JSON genérico em texto legível (chave: valor, recursivo)."""
    dados = json.loads(caminho.read_text(encoding="utf-8"))

    def _achatar(obj, prefixo=""):
        linhas = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                linhas.extend(_achatar(v, f"{prefixo}{k}."))
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                linhas.extend(_achatar(v, f"{prefixo}{i}."))
        else:
            linhas.append(f"{prefixo.rstrip('.')}: {obj}")
        return linhas

    return "\n".join(_achatar(dados))


def carregar_html(caminho: Path) -> str:
    """Remove tags HTML mantendo o texto visível. Requer beautifulsoup4."""
    try:
        from bs4 import BeautifulSoup
    except ImportError as e:
        raise ImportError(
            "beautifulsoup4 é necessário para carregar HTML: pip install beautifulsoup4"
        ) from e
    soup = BeautifulSoup(caminho.read_text(encoding="utf-8"), "html.parser")
    return soup.get_text(separator="\n", strip=True)


def carregar_pdf(caminho: Path) -> str:
    """Extrai texto de PDF nativo. Para PDFs escaneados, seria necessário OCR (não incluso aqui)."""
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise ImportError("pypdf é necessário para carregar PDF: pip install pypdf") from e
    reader = PdfReader(str(caminho))
    paginas = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(paginas).strip()


def carregar_docx(caminho: Path) -> str:
    """Extrai texto corrido de um .docx, preservando parágrafos."""
    try:
        import docx
    except ImportError as e:
        raise ImportError("python-docx é necessário para carregar Word: pip install python-docx") from e
    documento = docx.Document(str(caminho))
    return "\n".join(p.text for p in documento.paragraphs if p.text.strip())


def carregar_xlsx(caminho: Path) -> str:
    """Converte planilhas em texto estruturado, linha a linha com cabeçalhos repetidos."""
    try:
        import openpyxl
    except ImportError as e:
        raise ImportError("openpyxl é necessário para carregar Excel: pip install openpyxl") from e
    wb = openpyxl.load_workbook(str(caminho), data_only=True)
    blocos = []
    for aba in wb.worksheets:
        linhas = list(aba.iter_rows(values_only=True))
        if not linhas:
            continue
        cabecalho = linhas[0]
        for linha in linhas[1:]:
            partes = [f"{c}: {v}" for c, v in zip(cabecalho, linha) if v is not None]
            if partes:
                blocos.append(f"[{aba.title}] " + "; ".join(partes))
    return "\n".join(blocos)


def carregar_pptx(caminho: Path) -> str:
    """Extrai texto de cada slide, incluindo notas do apresentador."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError("python-pptx é necessário para carregar PowerPoint: pip install python-pptx") from e
    prs = Presentation(str(caminho))
    blocos = []
    for i, slide in enumerate(prs.slides, start=1):
        textos = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            textos.append(f"[Notas] {slide.notes_slide.notes_text_frame.text}")
        if textos:
            blocos.append(f"[Slide {i}] " + " | ".join(textos))
    return "\n".join(blocos)


_LOADERS = {
    ".md": carregar_markdown,
    ".markdown": carregar_markdown,
    ".csv": carregar_csv,
    ".json": carregar_json,
    ".html": carregar_html,
    ".htm": carregar_html,
    ".pdf": carregar_pdf,
    ".docx": carregar_docx,
    ".xlsx": carregar_xlsx,
    ".pptx": carregar_pptx,
}


def carregar_documento(caminho: Path) -> str:
    """Ponto de entrada único: detecta o formato pela extensão e chama o loader certo."""
    extensao = caminho.suffix.lower()
    loader = _LOADERS.get(extensao)
    if loader is None:
        raise LoaderNaoSuportadoError(f"Formato não suportado: {extensao} ({caminho})")
    return loader(caminho)
